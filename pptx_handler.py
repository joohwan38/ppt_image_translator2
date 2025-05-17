# pptx_handler.py
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID

import os
import io
import logging
import re
from datetime import datetime
import hashlib
import traceback
from PIL import Image
import tempfile
import shutil

import config
from interfaces import AbsPptxProcessor, AbsTranslator, AbsOcrHandler, AbsOllamaService
from utils import setup_task_logging # 3단계: 공통 로그 설정 함수 import

from typing import Optional, Dict, Any, List, Tuple, Callable, TypedDict, IO
from concurrent.futures import ThreadPoolExecutor


logger = logging.getLogger(__name__)

MEANINGFUL_CHAR_PATTERN = re.compile(
    r'[a-zA-Z'
    r'\u00C0-\u024F'    # Latin Extended-A
    r'\u1E00-\u1EFF'    # Latin Extended Additional
    r'\u0600-\u06FF'    # Arabic
    r'\u0750-\u077F'    # Arabic Supplement
    r'\u08A0-\u08FF'    # Arabic Extended-A
    r'\u3040-\u30ff'    # Hiragana, Katakana
    r'\u3131-\uD79D'    # Hangul Compatibility Jamo, Hangul Syllables
    r'\u4e00-\u9fff'    # CJK Unified Ideographs
    r'\u0E00-\u0E7F'    # Thai
    r']'
)

def should_skip_translation(text: str) -> bool:
    # (이전 코드와 동일)
    if not text: return True
    stripped_text = text.strip()
    if not stripped_text:
        logger.debug(f"번역 스킵 (공백만 존재): '{text[:50]}...'")
        return True
    
    if not MEANINGFUL_CHAR_PATTERN.search(stripped_text):
        logger.debug(f"번역 스킵 (의미 있는 문자 없음): '{stripped_text[:50]}...'")
        return True

    text_len = len(stripped_text)
    if text_len <= 3: 
        if MEANINGFUL_CHAR_PATTERN.search(stripped_text): 
            logger.debug(f"번역 시도 (짧은 문자열, 의미 문자 포함): '{stripped_text}'")
            return False
        else:
            logger.debug(f"번역 스킵 (짧고 의미 있는 문자 없음): '{stripped_text}'")
            return True

    meaningful_chars_count = len(MEANINGFUL_CHAR_PATTERN.findall(stripped_text))
    ratio = meaningful_chars_count / text_len
    if ratio < config.MIN_MEANINGFUL_CHAR_RATIO_SKIP:
        logger.debug(f"번역 스킵 (의미 문자 비율 낮음 {ratio:.2f}, 임계값: {config.MIN_MEANINGFUL_CHAR_RATIO_SKIP}): '{stripped_text[:50]}...'")
        return True

    logger.debug(f"번역 시도 (조건 통과): '{stripped_text[:50]}...'")
    return False


def is_ocr_text_valid(text: str) -> bool:
    # (이전 코드와 동일)
    if not text: return False
    stripped_text = text.strip()
    if not stripped_text: return False

    if not MEANINGFUL_CHAR_PATTERN.search(stripped_text):
        logger.debug(f"OCR 유효성 스킵 (의미 있는 문자 없음): '{stripped_text[:50]}...'")
        return False
    
    text_len = len(stripped_text)
    if text_len <= 2: 
        if MEANINGFUL_CHAR_PATTERN.search(stripped_text):
             logger.debug(f"OCR 유효 (매우 짧은 문자열, 의미 문자 포함): '{stripped_text}'")
             return True
        else:
            logger.debug(f"OCR 유효성 스킵 (매우 짧고 의미 있는 문자 없음): '{stripped_text}'")
            return False

    meaningful_chars_count = len(MEANINGFUL_CHAR_PATTERN.findall(stripped_text))
    ratio = meaningful_chars_count / text_len
    if ratio < config.MIN_MEANINGFUL_CHAR_RATIO_OCR:
        logger.debug(f"OCR 유효성 스킵 (의미 문자 비율 낮음 {ratio:.2f}, 임계값: {config.MIN_MEANINGFUL_CHAR_RATIO_OCR}): '{stripped_text[:50]}...'")
        return False
        
    logger.debug(f"OCR 유효 (조건 통과): '{stripped_text[:50]}...'")
    return True

class TranslationJob(TypedDict):
    original_text: str
    context: Dict[str, Any]
    is_ocr: bool
    char_count: int

class PptxHandler(AbsPptxProcessor):
    def __init__(self):
        pass

    def get_file_info(self, file_path: str) -> Dict[str, int]:
        # (이전 코드와 동일)
        logger.info(f"파일 정보 분석 시작: {file_path}")
        info = {
            "slide_count": 0,
            "text_elements_count": 0,
            "total_text_char_count": 0,
            "image_elements_count": 0,
            "chart_elements_count": 0
        }
        try:
            prs = Presentation(file_path)
            info["slide_count"] = len(prs.slides)
            for slide in prs.slides:
                for shape in slide.shapes:
                    is_counted_as_text_element_for_shape = False
                    if shape.has_text_frame and hasattr(shape.text_frame, 'text') and \
                       shape.text_frame.text and shape.text_frame.text.strip():
                        text_content = shape.text_frame.text
                        if not should_skip_translation(text_content):
                            if not is_counted_as_text_element_for_shape:
                                info["text_elements_count"] += 1
                                is_counted_as_text_element_for_shape = True
                            info["total_text_char_count"] += len(text_content)

                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        info["image_elements_count"] += 1
                    elif shape.has_table:
                        for r_idx, row in enumerate(shape.table.rows):
                            for c_idx, cell in enumerate(row.cells):
                                if hasattr(cell.text_frame, 'text') and cell.text_frame.text and cell.text_frame.text.strip():
                                    text_content = cell.text_frame.text
                                    if not should_skip_translation(text_content):
                                        info["text_elements_count"] += 1
                                        info["total_text_char_count"] += len(text_content)
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        info["chart_elements_count"] +=1
            logger.info(
                f"파일 분석 완료: Slides:{info['slide_count']}, "
                f"TextElements(internal):{info['text_elements_count']} (TotalChars:{info['total_text_char_count']}), "
                f"Images:{info['image_elements_count']}, Charts:{info['chart_elements_count']}"
            )
        except Exception as e:
            logger.error(f"'{os.path.basename(file_path)}' 파일 정보 분석 오류: {e}", exc_info=True)
            # 오류 시 기본값 반환 또는 특정 오류 값 설정 가능
            info = { "slide_count": -1, "text_elements_count": -1, "total_text_char_count": -1,
                     "image_elements_count": -1, "chart_elements_count": -1, "error": str(e)}
        return info


    def _get_style_properties(self, font_object) -> Dict[str, Any]:
        # (이전 코드와 동일)
        if font_object is None:
            return {}
        style_props: Dict[str, Any] = {
            'name': None, 'size': None, 'bold': None, 'italic': None,
            'underline': None, 'color_rgb': None, 'color_theme_index': None,
            'color_brightness': 0.0, 'language_id': None
        }
        try: style_props['name'] = font_object.name
        except AttributeError: pass
        try: style_props['size'] = font_object.size
        except AttributeError: pass
        try: style_props['bold'] = font_object.bold
        except AttributeError: pass
        try: style_props['italic'] = font_object.italic
        except AttributeError: pass
        try: style_props['underline'] = font_object.underline
        except AttributeError: pass
        try: style_props['language_id'] = font_object.language_id
        except (ValueError, AttributeError): pass

        if hasattr(font_object, 'color') and hasattr(font_object.color, 'type'):
            color_type = font_object.color.type
            if color_type == MSO_COLOR_TYPE.RGB:
                try: style_props['color_rgb'] = tuple(font_object.color.rgb)
                except AttributeError: pass
            elif color_type == MSO_COLOR_TYPE.SCHEME:
                try: style_props['color_theme_index'] = font_object.color.theme_color
                except AttributeError: pass
                try: style_props['color_brightness'] = font_object.color.brightness
                except AttributeError: style_props['color_brightness'] = 0.0
        return style_props


    def _get_text_style(self, run) -> Dict[str, Any]:
        # (이전 코드와 동일)
        style = self._get_style_properties(run.font)
        try:
            style['hyperlink_address'] = run.hyperlink.address if run.hyperlink and run.hyperlink.address else None
        except AttributeError:
            style['hyperlink_address'] = None
        return style

    def _apply_style_properties(self, target_font_object, style_dict_to_apply: Dict[str, Any]):
        # (이전 코드와 동일)
        if not style_dict_to_apply or target_font_object is None: return
        font = target_font_object
        if style_dict_to_apply.get('name') is not None: font.name = style_dict_to_apply['name']
        if style_dict_to_apply.get('size') is not None: font.size = style_dict_to_apply['size']
        if style_dict_to_apply.get('bold') is not None: font.bold = style_dict_to_apply['bold']
        if style_dict_to_apply.get('italic') is not None: font.italic = style_dict_to_apply['italic']
        if style_dict_to_apply.get('underline') is not None: font.underline = style_dict_to_apply['underline']
        applied_color = False
        if style_dict_to_apply.get('color_rgb') is not None:
            try:
                rgb_tuple = tuple(int(c) for c in style_dict_to_apply['color_rgb'])
                if len(rgb_tuple) == 3:
                    font.color.rgb = RGBColor(*rgb_tuple)
                    applied_color = True
            except Exception as e: logger.warning(f"RGB 색상 {style_dict_to_apply['color_rgb']} 적용 실패: {e}")
        if not applied_color and style_dict_to_apply.get('color_theme_index') is not None:
            try:
                theme_color_val = style_dict_to_apply['color_theme_index']
                if isinstance(theme_color_val, MSO_THEME_COLOR_INDEX):
                    font.color.theme_color = theme_color_val
                elif isinstance(theme_color_val, int) and theme_color_val in [item.value for item in MSO_THEME_COLOR_INDEX]:
                    font.color.theme_color = MSO_THEME_COLOR_INDEX(theme_color_val)
                else: logger.warning(f"유효하지 않은 테마 색상 인덱스: {theme_color_val}")
                brightness_val = style_dict_to_apply.get('color_brightness', 0.0)
                font.color.brightness = max(-1.0, min(1.0, float(brightness_val)))
            except Exception as e: logger.warning(f"테마 색상 {style_dict_to_apply.get('color_theme_index')} 적용 실패: {e}")
        if style_dict_to_apply.get('language_id') is not None:
            try:
                lang_id_val = style_dict_to_apply['language_id']
                if isinstance(lang_id_val, MSO_LANGUAGE_ID):
                    font.language_id = lang_id_val
                elif isinstance(lang_id_val, int):
                    try: font.language_id = MSO_LANGUAGE_ID(lang_id_val)
                    except ValueError: logger.warning(f"MSO_LANGUAGE_ID에 없는 언어 ID: {lang_id_val}. 기본값 사용.")
                else: logger.debug(f"문자열 언어 ID '{lang_id_val}'는 직접 적용 불가. 무시됨.")
            except Exception as e_lang: logger.warning(f"language_id '{style_dict_to_apply['language_id']}' 적용 실패: {e_lang}")

    def _apply_text_style(self, run, style_to_apply: Dict[str, Any]):
        # (이전 코드와 동일)
        if not style_to_apply or run is None: return
        self._apply_style_properties(run.font, style_to_apply)
        # 하이퍼링크 적용 로직 (필요 시 추가)


    def translate_presentation_stage1(self, prs: Presentation, src_lang_ui_name: str, tgt_lang_ui_name: str,
                                      translator: AbsTranslator, ocr_handler: Optional[AbsOcrHandler],
                                      model_name: str, ollama_service: AbsOllamaService,
                                      font_code_for_render: str, task_log_filepath: str,
                                      progress_callback_item_completed: Optional[Callable[[Any, str, int, str], None]] = None,
                                      stop_event: Optional[Any] = None,
                                      image_translation_enabled: bool = True,
                                      ocr_temperature: Optional[float] = None
                                      ) -> bool:

        # --- 3단계: 공통 로그 파일 생성 로직 사용 ---
        initial_log_lines_s1 = ["--- 1단계: 차트 외 요소 번역 시작 (translate_presentation_stage1) ---"]
        f_task_log_s1: Optional[IO[str]] = None
        log_func_s1: Optional[Callable[[str], None]] = None

        if task_log_filepath:
            f_task_log_s1, log_func_s1_temp = setup_task_logging(task_log_filepath, initial_log_lines_s1)
            if log_func_s1_temp: log_func_s1 = log_func_s1_temp
        # log_func_s1이 None일 경우, 주요 로깅은 logger.info/warning/error를 통해 이루어짐

        logger.info("1단계: 차트 외 요소 (텍스트 상자, 표, OCR 등) 수집 중...")
        if log_func_s1: log_func_s1("1단계: 차트 외 요소 수집 중...")


        translation_jobs: List[TranslationJob] = []
        # image_update_jobs는 현재 직접 사용되지 않으므로 제거하거나 주석 처리 가능
        # image_update_jobs: List[Dict[str, Any]] = []
        elements_to_analyze_stage1: List[Dict[str, Any]] = [] # UI 진행 표시용
        original_paragraph_styles_stage1: Dict[Tuple[int, Any, Any], List[Dict[str, Any]]] = {}


        for slide_idx, slide in enumerate(prs.slides):
            if stop_event and stop_event.is_set(): break
            for shape_idx, shape in enumerate(slide.shapes):
                if stop_event and stop_event.is_set(): break

                shape_id = getattr(shape, 'shape_id', f"slide{slide_idx}_shape{shape_idx}")
                element_name = shape.name or f"S{slide_idx+1}_Shape{shape_idx}_Id{shape_id}"
                item_base_info = {'slide_idx': slide_idx, 'shape_obj_ref': shape, 'name': element_name, 'shape_id_log': shape_id}
                item_progress_type = "Unknown" # UI 피드백용 타입 문자열

                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    elements_to_analyze_stage1.append({**item_base_info, 'type': 'chart_placeholder', 'char_count':0, 'progress_type': "차트 (2단계 처리)"})
                    continue

                if shape.has_text_frame and hasattr(shape.text_frame, 'text') and \
                   shape.text_frame.text and shape.text_frame.text.strip():
                    original_text = shape.text_frame.text
                    char_count = 0
                    if not should_skip_translation(original_text):
                        char_count = len(original_text)
                        style_key_suffix_val = 'shape_text'
                        # id(shape) 대신 shape.shape_id 사용 시도, 없으면 기존 방식
                        style_unique_key_id_part = shape_id if hasattr(shape, 'shape_id') else id(shape)
                        style_unique_key = (slide_idx, style_unique_key_id_part, style_key_suffix_val)
                        translation_jobs.append({
                            'original_text': original_text,
                            'context': {**item_base_info, 'item_type_internal': 'text_shape', 'style_unique_key': style_unique_key},
                            'is_ocr': False,
                            'char_count': char_count
                        })
                    item_progress_type = "텍스트 요소" # --- 3단계: UI 피드백 구체화 ---
                    elements_to_analyze_stage1.append({**item_base_info, 'type': 'text_shape', 'original_text': original_text, 'char_count': char_count, 'progress_type': item_progress_type})

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    item_progress_type = "이미지 (OCR 대상)" # --- 3단계: UI 피드백 구체화 ---
                    elements_to_analyze_stage1.append({**item_base_info, 'type': 'image', 'progress_type': item_progress_type, 'char_count':0})
                    # OCR 작업은 아래에서 일괄 처리

                elif shape.has_table:
                    item_progress_type = "표 내부 텍스트" # --- 3단계: UI 피드백 구체화 ---
                    elements_to_analyze_stage1.append({**item_base_info, 'type': 'table_container', 'progress_type': item_progress_type, 'char_count':0})
                    for r_idx, row in enumerate(shape.table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            if hasattr(cell.text_frame, 'text') and cell.text_frame.text and cell.text_frame.text.strip():
                                original_text = cell.text_frame.text
                                char_count = 0
                                if not should_skip_translation(original_text):
                                    char_count = len(original_text)
                                    style_key_suffix_val = (r_idx, c_idx)
                                    style_unique_key_id_part = shape_id if hasattr(shape, 'shape_id') else id(shape)
                                    style_unique_key = (slide_idx, style_unique_key_id_part, style_key_suffix_val)
                                    cell_item_name = f"{element_name}_R{r_idx}C{c_idx}"
                                    translation_jobs.append({
                                        'original_text': original_text,
                                        'context': {
                                            'slide_idx': slide_idx, 'table_shape_obj_ref': shape,
                                            'name': cell_item_name, 'item_type_internal': 'table_cell',
                                            'row_idx': r_idx, 'col_idx': c_idx, 'style_unique_key': style_unique_key,
                                            'shape_id_log': shape_id
                                        },
                                        'is_ocr': False,
                                        'char_count': char_count
                                    })
        if log_func_s1:
            log_func_s1(f"1단계 분석 대상 요소 (UI 진행 표시용): {len(elements_to_analyze_stage1)}개.")
            log_func_s1(f"1단계 번역 작업 수집 완료 (텍스트, 표): {len(translation_jobs)}개 항목.")
        logger.info(f"1단계 번역 작업 수집 완료 (텍스트, 표): {len(translation_jobs)}개 항목.")

        if not translation_jobs and not (image_translation_enabled and ocr_handler):
             is_any_chart_present = any(s.shape_type == MSO_SHAPE_TYPE.CHART for slide_obj in prs.slides for s in slide_obj.shapes)
             if not is_any_chart_present:
                msg = "1단계 번역/처리 대상 요소(텍스트/표)가 없고, OCR 비활성화 또는 핸들러 부재, 차트도 없어 1단계 처리 스킵."
                if log_func_s1: log_func_s1(msg)
                logger.info(msg)
             if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
             return True

        texts_for_batch_translation = [job['original_text'] for job in translation_jobs if job['char_count'] > 0 and not job['is_ocr']]
        translated_texts_batch: List[str] = []
        if texts_for_batch_translation:
            if log_func_s1: log_func_s1(f"일반 텍스트 {len(texts_for_batch_translation)}개 일괄 번역 시작...")
            logger.info(f"일반 텍스트 {len(texts_for_batch_translation)}개 일괄 번역 시작...")
            translated_texts_batch = translator.translate_texts_batch(
                texts_for_batch_translation, src_lang_ui_name, tgt_lang_ui_name,
                model_name, ollama_service, is_ocr_text=False, stop_event=stop_event
            )
            if log_func_s1: log_func_s1(f"일반 텍스트 일괄 번역 완료. 결과 {len(translated_texts_batch)}개 받음.")
            logger.info(f"일반 텍스트 일괄 번역 완료. 결과 {len(translated_texts_batch)}개 받음.")
            if len(texts_for_batch_translation) != len(translated_texts_batch):
                warn_msg = f"경고: 원본 텍스트 수({len(texts_for_batch_translation)})와 번역 결과 수({len(translated_texts_batch)}) 불일치!"
                if log_func_s1: log_func_s1(warn_msg)
                logger.warning(warn_msg + " 배치 번역 로직 점검 필요.")
                if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                return False

        translated_text_idx = 0
        for job_idx, job_data in enumerate(translation_jobs):
            # ... (이하 텍스트/표 번역 적용 로직은 이전과 거의 동일하게 유지, log_func_s1 사용) ...
            # --- 3단계: UI 피드백 구체화 (progress_callback_item_completed 호출 시 전달 정보) ---
            # 예시: current_task_type을 "텍스트 요소 번역 적용", "표 셀 번역 적용" 등으로 구체화
            if stop_event and stop_event.is_set():
                if log_func_s1: log_func_s1(f"1단계 적용 중 중단 요청 감지.")
                if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
                return False

            context = job_data['context']
            slide_idx = context['slide_idx']
            item_name_log = context['name']
            item_type_internal = context['item_type_internal']
            shape_id_log = context['shape_id_log']
            current_progress_text = job_data['original_text'][:30] # UI 표시용 텍스트 조각
            weighted_work_for_item = job_data['char_count'] * config.WEIGHT_TEXT_CHAR
            
            # UI 피드백용 작업 타입 문자열 결정
            ui_task_type_feedback = "텍스트/표 적용" # 기본값
            if item_type_internal == 'text_shape':
                ui_task_type_feedback = "텍스트 요소 적용"
            elif item_type_internal == 'table_cell':
                ui_task_type_feedback = "표 셀 내용 적용"


            if log_func_s1: log_func_s1(f"  [1단계 S{slide_idx+1}] 적용 시작: '{item_name_log}' (ID: {shape_id_log}), 타입: {item_type_internal}")

            text_frame_for_processing: Optional[Any] = None
            # ... (text_frame_for_processing 할당 로직) ...
            if item_type_internal == 'text_shape':
                shape_obj = context.get('shape_obj_ref')
                if shape_obj and shape_obj.has_text_frame:
                    text_frame_for_processing = shape_obj.text_frame
            elif item_type_internal == 'table_cell':
                table_shape_obj = context.get('table_shape_obj_ref')
                row_idx, col_idx = context['row_idx'], context['col_idx']
                if table_shape_obj and table_shape_obj.has_table:
                    try: text_frame_for_processing = table_shape_obj.table.cell(row_idx, col_idx).text_frame
                    except IndexError:
                        err_msg_tbl_idx = f"1단계 테이블 셀 접근 오류 (IndexError): {item_name_log} at R{row_idx}C{col_idx}"
                        logger.error(err_msg_tbl_idx);
                        if log_func_s1: log_func_s1(f"    오류: {err_msg_tbl_idx}. 건너뜀.")
                        if progress_callback_item_completed: progress_callback_item_completed(f"슬라이드 {slide_idx + 1} 오류", "테이블 접근 실패", weighted_work_for_item, f"셀 ({row_idx},{col_idx})")
                        continue


            if text_frame_for_processing and job_data['char_count'] > 0:
                style_unique_key = context['style_unique_key']
                if style_unique_key not in original_paragraph_styles_stage1:
                    para_styles_collected: List[Dict[str, Any]] = []
                    for para_obj in text_frame_for_processing.paragraphs:
                        para_default_font_style = self._get_style_properties(para_obj.font)
                        runs_info: List[Dict[str, Any]] = []
                        if para_obj.runs:
                            for run_obj in para_obj.runs:
                                runs_info.append({'text': run_obj.text, 'style': self._get_text_style(run_obj)})
                        elif para_obj.text and para_obj.text.strip():
                            run_style_from_para = para_default_font_style.copy()
                            run_style_from_para['hyperlink_address'] = None
                            runs_info.append({'text': para_obj.text, 'style': run_style_from_para})
                        para_styles_collected.append({
                            'runs': runs_info,
                            'alignment': para_obj.alignment,
                            'level': para_obj.level,
                            'space_before': para_obj.space_before,
                            'space_after': para_obj.space_after,
                            'line_spacing': para_obj.line_spacing,
                            'paragraph_default_run_style': para_default_font_style
                        })
                    original_paragraph_styles_stage1[style_unique_key] = para_styles_collected
                    if log_func_s1:
                        log_func_s1(f"      '{item_name_log}'의 원본 단락 스타일 저장 (1단계 적용 시점).")

                translated_text_content = translated_texts_batch[translated_text_idx] if translated_text_idx < len(translated_texts_batch) else job_data['original_text']
                translated_text_idx += 1
                current_progress_text = translated_text_content[:30].replace('\n', ' ')
                log_trans_text_snippet = translated_text_content.replace('\n', ' / ').strip()[:100]
                
                if log_func_s1:
                    log_func_s1(f"    [1단계 적용 전] \"{job_data['original_text'].strip()[:50]}...\" -> [1단계 적용 후] \"{log_trans_text_snippet}...\"")

                if "오류:" not in translated_text_content: # 번역 성공 시에만 적용 (오류 메시지는 텍스트로 삽입하지 않음)
                    stored_paras_info_apply = original_paragraph_styles_stage1.get(style_unique_key, [])
                    
                    original_tf_auto_sz = getattr(text_frame_for_processing, 'auto_size', None)
                    original_tf_vertical_anchor = getattr(text_frame_for_processing, 'vertical_anchor', None)
                    original_tf_word_wrap = getattr(text_frame_for_processing, 'word_wrap', None)
                    original_tf_margin_left = getattr(text_frame_for_processing, 'margin_left', None)
                    original_tf_margin_right = getattr(text_frame_for_processing, 'margin_right', None)
                    original_tf_margin_top = getattr(text_frame_for_processing, 'margin_top', None)
                    original_tf_margin_bottom = getattr(text_frame_for_processing, 'margin_bottom', None)

                    if original_tf_auto_sz is not None and original_tf_auto_sz != MSO_AUTO_SIZE.NONE:
                        try:
                            text_frame_for_processing.auto_size = MSO_AUTO_SIZE.NONE
                        except Exception as e_auto_sz_none:
                            logger.debug(f"auto_size=NONE 설정 중 예외 (무시): {e_auto_sz_none}")

                    if original_tf_word_wrap is not None and not original_tf_word_wrap: # 워드랩이 False였으면 True로 변경 (다국어 줄바꿈 문제 방지)
                        try:
                            text_frame_for_processing.word_wrap = True
                        except Exception as e_word_wrap_true:
                            logger.debug(f"word_wrap=True 설정 중 예외 (무시): {e_word_wrap_true}")
                    
                    text_frame_for_processing.clear() # 기존 단락들 제거

                    # --- 사용자 요청 코드 시작 ---
                    if not translated_text_content.strip(): # 번역 결과가 공백/빈 문자열인 경우
                        if log_func_s1: log_func_s1(f"      번역 결과가 비어 공백 단락 1개 추가 요청됨 (내용: '{translated_text_content}').")
                        p_new = text_frame_for_processing.add_paragraph()
                        run_new = p_new.add_run()
                        run_new.text = " " # 빈 공간이라도 단락 유지
                        if stored_paras_info_apply: # 첫번째 단락 스타일이라도 적용
                            para_template_empty = stored_paras_info_apply[0]
                            if para_template_empty.get('alignment'): p_new.alignment = para_template_empty['alignment']
                            p_new.level = para_template_empty.get('level', 0)
                            p_new.space_before = para_template_empty.get('space_before', Pt(0))
                            p_new.space_after = para_template_empty.get('space_after', Pt(0))
                            p_new.line_spacing = para_template_empty.get('line_spacing')
                            if para_template_empty.get('runs') and para_template_empty['runs'][0].get('style'):
                                    self._apply_text_style(run_new, para_template_empty['runs'][0]['style'])
                        if log_func_s1: log_func_s1(f"        번역 결과가 비어 공백 단락 1개 추가 완료.")
                    else:
                        # 실제 내용이 있을 때는 다른 방식으로 접근
                        # text_frame을 clear 후 XML 하위 요소 확인
                        if hasattr(text_frame_for_processing, '_element') and text_frame_for_processing._element is not None:
                            # 기존에 생성된 빈 단락을 모두 제거 (XML 레벨에서)
                            txBody_xml_element = text_frame_for_processing._element # txBody
                            child_elements_to_remove = [child for child in list(txBody_xml_element) if child.tag.endswith('}p')] # 'a:p' 태그
                            if child_elements_to_remove:
                                if log_func_s1: log_func_s1(f"        XML 레벨에서 기존 단락 {len(child_elements_to_remove)}개 제거 시도...")
                                for child_xml_p_tag in child_elements_to_remove:
                                    txBody_xml_element.remove(child_xml_p_tag)
                                if log_func_s1: log_func_s1(f"        XML 레벨에서 기존 단락 제거 완료.")
                            else:
                                if log_func_s1: log_func_s1(f"        XML 레벨에서 제거할 기존 단락 없음.")
                        
                        # 이제 단락 추가
                        lines = translated_text_content.splitlines()
                        if not lines and translated_text_content: # splitlines 결과가 비었지만 원본은 내용이 있는 경우 (예: "내용" -> ["내용"])
                            lines = [translated_text_content]
                        elif not lines: # 그래도 비어있으면 공백 한 줄로 처리
                            lines = [" "]
                        
                        for i, line in enumerate(lines):
                            p = text_frame_for_processing.add_paragraph()
                            # 스타일 적용
                            para_template = stored_paras_info_apply[min(i, len(stored_paras_info_apply)-1)] if stored_paras_info_apply else {}
                            if para_template:
                                if para_template.get('alignment') is not None: p.alignment = para_template['alignment']
                                p.level = para_template.get('level', 0)
                                # 모든 단락에 동일하게 적용 (이전에는 첫 단락에만 적용되던 문제를 수정)
                                p.space_before = para_template.get('space_before', Pt(0)) # 명시적으로 Pt(0) 또는 저장된 값
                                p.space_after = para_template.get('space_after', Pt(0)) # 명시적으로 Pt(0) 또는 저장된 값
                                p.line_spacing = para_template.get('line_spacing')
                            
                            run = p.add_run()
                            run.text = line if line.strip() else " " # 빈 줄 대신 공백 한 칸 삽입 (단락 유지)
                            
                            if para_template and para_template.get('runs'):
                                # 첫 번째 run의 스타일을 현재 run에 적용 (가장 기본적인 스타일 복원 방식)
                                # 더 정교하게 하려면 원본 run 개수와 매칭하거나, 줄바꿈 기준으로 나눠진 텍스트에 각각 스타일 적용 필요
                                if para_template['runs']: # 비어있지 않은지 확인
                                    original_run_style_to_apply = para_template['runs'][0]['style']
                                    self._apply_text_style(run, original_run_style_to_apply)
                            elif para_template and para_template.get('paragraph_default_run_style'): # 단락 기본 스타일이라도 적용
                                self._apply_text_style(run, para_template['paragraph_default_run_style'])

                        if log_func_s1: log_func_s1(f"        '{item_name_log}'에 번역된 텍스트 ({len(lines)}줄) 및 스타일 적용 완료.")
                    # --- 사용자 요청 코드 종료 ---

                    # 원래 텍스트 프레임 속성 복원
                    if original_tf_auto_sz is not None:
                        try: text_frame_for_processing.auto_size = original_tf_auto_sz
                        except Exception as e_auto_sz_restore: logger.debug(f"auto_size 복원 중 예외 (무시): {e_auto_sz_restore}")
                    if original_tf_vertical_anchor is not None:
                        try: text_frame_for_processing.vertical_anchor = original_tf_vertical_anchor
                        except Exception as e_va_restore: logger.debug(f"vertical_anchor 복원 중 예외 (무시): {e_va_restore}")
                    
                    # word_wrap은 True로 유지하거나, 원래 False였더라도 True로 두는 것이 다국어 환경에 유리할 수 있음
                    # 여기서는 원래 False였더라도 True로 유지 (위에서 이미 설정)
                    # if original_tf_word_wrap is not None:
                    #     try: text_frame_for_processing.word_wrap = original_tf_word_wrap
                    #     except Exception as e_ww_restore: logger.debug(f"word_wrap 복원 중 예외 (무시): {e_ww_restore}")

                    if original_tf_margin_left is not None: text_frame_for_processing.margin_left = original_tf_margin_left
                    if original_tf_margin_right is not None: text_frame_for_processing.margin_right = original_tf_margin_right
                    if original_tf_margin_top is not None: text_frame_for_processing.margin_top = original_tf_margin_top
                    if original_tf_margin_bottom is not None: text_frame_for_processing.margin_bottom = original_tf_margin_bottom
                    
                    if log_func_s1: log_func_s1(f"        '{item_name_log}' 1단계 번역된 텍스트 적용 완료 (최종).")
                else: # 번역 실패 또는 빈 결과 (오류 메시지 포함)
                    if log_func_s1:
                        log_func_s1(f"      -> 1단계 텍스트 번역 실패 또는 빈 결과로 내용 변경 없음: \"{translated_text_content.strip()[:50]}...\"")
            
            elif job_data['char_count'] == 0 and item_type_internal in ['text_shape', 'table_cell']:
                if log_func_s1:
                    log_func_s1(f"      [1단계 스킵됨 - 번역 불필요 또는 의미 없는 텍스트]")


            if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", ui_task_type_feedback, weighted_work_for_item, current_progress_text)
            if log_func_s1: log_func_s1("\n") # 각 job 처리 후 빈 줄 추가 (로그 가독성)


        # 이미지 OCR 및 번역 적용 로직
        if image_translation_enabled and ocr_handler:
            # --- 3단계: OCR 처리 중 더 세분화된 피드백 (progress_callback_item_completed 호출 시 전달 정보) ---
            for slide_idx, slide_obj in enumerate(prs.slides):
                if stop_event and stop_event.is_set(): break
                shapes_to_process_for_ocr = [s for s in slide_obj.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
                for shape_obj_ocr in shapes_to_process_for_ocr:
                    if stop_event and stop_event.is_set(): break
                    shape_id_ocr = getattr(shape_obj_ocr, 'shape_id', f"slide{slide_idx}_shape_ocr_{id(shape_obj_ocr)}")
                    item_name_ocr = shape_obj_ocr.name or f"S{slide_idx+1}_ImgOCR_Id{shape_id_ocr}"
                    if log_func_s1: log_func_s1(f"  [1단계 S{slide_idx+1}] OCR 처리 시도: '{item_name_ocr}'")
                    weighted_work_for_ocr_item = config.WEIGHT_IMAGE
                    current_ocr_progress_text = f"이미지 '{item_name_ocr}' OCR 중..." # 초기 피드백

                    # OCR 진행 콜백 업데이트 (OCR 시작)
                    if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                        progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", "이미지 OCR 시작", 0, f"'{item_name_ocr}' 분석 중")


                    try:
                        img_bytes_io = io.BytesIO(shape_obj_ocr.image.blob)
                        with Image.open(img_bytes_io) as img_pil_original_ocr:
                            img_pil_rgb_ocr = img_pil_original_ocr.convert("RGB")
                            ocr_results_list = ocr_handler.ocr_image(img_pil_rgb_ocr) # 여기서 시간 소요

                            # OCR 결과 후 피드백
                            current_ocr_progress_text = f"'{item_name_ocr}' OCR 완료, {len(ocr_results_list) if ocr_results_list else 0}개 텍스트 블록"
                            if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", "이미지 OCR 완료", 0, current_ocr_progress_text)


                            if ocr_results_list:
                                # ... (OCR 결과 처리, 번역, 렌더링 로직) ...
                                # 각 텍스트 블록 번역 시 또는 렌더링 시 추가 피드백 가능
                                # 예: progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", "이미지 내 텍스트 번역 중", ..., f"'{item_name_ocr}'의 텍스트 블록 X 번역")
                                # 예: progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", "이미지 내 텍스트 렌더링 중", ..., f"'{item_name_ocr}'의 텍스트 블록 X 렌더링")
                                if log_func_s1: log_func_s1(f"        이미지 내 OCR 텍스트 {len(ocr_results_list)}개 블록 발견.")
                                ocr_texts_to_translate_current_image = []
                                ocr_job_contexts_current_image = []
                                for ocr_idx, ocr_res_item in enumerate(ocr_results_list):
                                    # ... (유효성 검사 및 번역 대상 수집) ...
                                    if not (isinstance(ocr_res_item, (list, tuple)) and len(ocr_res_item) >= 2): continue
                                    ocr_box_coords, ocr_text_conf_pair = ocr_res_item[0], ocr_res_item[1]
                                    ocr_angle_info = ocr_res_item[2] if len(ocr_res_item) > 2 else None
                                    if not (isinstance(ocr_text_conf_pair, (list, tuple)) and len(ocr_text_conf_pair) == 2): continue
                                    ocr_text_original, ocr_confidence = ocr_text_conf_pair
                                    if is_ocr_text_valid(ocr_text_original) and not should_skip_translation(ocr_text_original):
                                        ocr_texts_to_translate_current_image.append(ocr_text_original)
                                        ocr_job_contexts_current_image.append({'box': ocr_box_coords, 'original_text': ocr_text_original, 'angle': ocr_angle_info, 'confidence': ocr_confidence})
                                    else:
                                        if log_func_s1: log_func_s1(f"          OCR Text 스킵됨 (유효성/번역 불필요): \"{ocr_text_original.strip()[:30]}...\"")

                                if ocr_texts_to_translate_current_image:
                                    if log_func_s1: log_func_s1(f"        이미지 내 유효 OCR 텍스트 {len(ocr_texts_to_translate_current_image)}개 배치 번역 시작...")
                                    # 번역 시작 전 피드백
                                    current_ocr_progress_text = f"'{item_name_ocr}' 텍스트 {len(ocr_texts_to_translate_current_image)}개 번역 중..."
                                    if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                         progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", "이미지 내 텍스트 번역", 0, current_ocr_progress_text)

                                    translated_ocr_texts_batch = translator.translate_texts_batch(ocr_texts_to_translate_current_image, src_lang_ui_name, tgt_lang_ui_name, model_name, ollama_service, is_ocr_text=True, ocr_temperature=ocr_temperature, stop_event=stop_event)
                                    
                                    current_ocr_progress_text = f"'{item_name_ocr}' 텍스트 번역 완료, 렌더링 중..."
                                    if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                                         progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", "이미지 내 텍스트 렌더링", 0, current_ocr_progress_text)


                                    if log_func_s1: log_func_s1(f"        이미지 내 OCR 텍스트 배치 번역 완료. 결과 {len(translated_ocr_texts_batch)}개 받음.")
                                    if len(ocr_texts_to_translate_current_image) == len(translated_ocr_texts_batch):
                                        img_bytes_io.seek(0) # 다시 사용하기 위해
                                        with Image.open(img_bytes_io) as img_to_render_on_base: # 원본 이미지 다시 열기
                                            original_img_format = img_to_render_on_base.format
                                            edited_img_pil = img_to_render_on_base.copy()
                                            any_ocr_text_rendered = False
                                            for i, translated_ocr_text_val in enumerate(translated_ocr_texts_batch):
                                                if stop_event and stop_event.is_set(): break
                                                ocr_job_ctx = ocr_job_contexts_current_image[i]
                                                # 개별 텍스트 렌더링 피드백 (선택적, 너무 많을 수 있음)
                                                # current_ocr_progress_text = f"'{item_name_ocr}'의 '{translated_ocr_text_val[:10]}...' 렌더링"
                                                # if progress_callback_item_completed: progress_callback_item_completed(f"슬라이드 {slide_idx + 1}", "이미지 텍스트 렌더링", 0, current_ocr_progress_text)

                                                if log_func_s1: log_func_s1(f"          OCR Text [{i+1}]: \"{ocr_job_ctx['original_text'].strip()[:30]}...\" -> 번역: \"{translated_ocr_text_val.strip()[:30]}...\"")
                                                if "오류:" not in translated_ocr_text_val and translated_ocr_text_val.strip():
                                                    try:
                                                        edited_img_pil = ocr_handler.render_translated_text_on_image(edited_img_pil, ocr_job_ctx['box'], translated_ocr_text_val, font_code_for_render=font_code_for_render, original_text=ocr_job_ctx['original_text'], ocr_angle=ocr_job_ctx['angle'])
                                                        any_ocr_text_rendered = True
                                                        if log_func_s1: log_func_s1(f"              -> 렌더링 완료.")
                                                    except Exception as e_render:
                                                        if log_func_s1: log_func_s1(f"              오류: OCR 텍스트 렌더링 실패: {e_render}")
                                                        logger.error(f"OCR 텍스트 렌더링 실패 ('{item_name_ocr}'): {e_render}", exc_info=True)
                                                else:
                                                    if log_func_s1: log_func_s1(f"            -> 번역 실패 또는 빈 결과로 렌더링 안 함.")
                                            if stop_event and stop_event.is_set(): break
                                            if any_ocr_text_rendered:
                                                # ... (이미지 교체 로직) ...
                                                output_img_stream = io.BytesIO()
                                                save_format_ocr_img = original_img_format if original_img_format and original_img_format.upper() in ['JPEG', 'PNG', 'GIF', 'BMP', 'TIFF'] else 'PNG'
                                                edited_img_pil.save(output_img_stream, format=save_format_ocr_img)
                                                output_img_stream.seek(0)
                                                left, top, width, height = shape_obj_ocr.left, shape_obj_ocr.top, shape_obj_ocr.width, shape_obj_ocr.height
                                                name_orig_img = shape_obj_ocr.name
                                                sp_xml_elem = shape_obj_ocr.element
                                                parent_xml_elem = sp_xml_elem.getparent()
                                                if parent_xml_elem is not None:
                                                    parent_xml_elem.remove(sp_xml_elem)
                                                    new_pic_shape = prs.slides[slide_idx].shapes.add_picture(output_img_stream, left, top, width=width, height=height)
                                                    if name_orig_img: new_pic_shape.name = name_orig_img
                                                    if log_func_s1: log_func_s1(f"        이미지 '{item_name_ocr}' 성공적으로 교체됨.")
                                                else:
                                                    if log_func_s1: log_func_s1(f"        경고: 이미지 '{item_name_ocr}'의 부모 XML 요소 찾지 못해 교체 실패. 원본 유지.")
                                            else:
                                                if log_func_s1: log_func_s1(f"        이미지 '{item_name_ocr}'에 번역 및 렌더링된 텍스트가 없어 변경 없음.")
                                    else:
                                        if log_func_s1: log_func_s1(f"        경고: 이미지 '{item_name_ocr}'의 OCR 텍스트 수와 번역 결과 수 불일치. 이미지 변경 없음.")
                                else:
                                    if log_func_s1: log_func_s1(f"        이미지 '{item_name_ocr}' 내 번역 대상 유효 OCR 텍스트 없음.")
                            else:
                                if log_func_s1: log_func_s1(f"        이미지 '{item_name_ocr}' 내에서 OCR 텍스트 발견되지 않음.")

                    except Exception as e_ocr_general_img:
                        err_msg_ocr_gen = f"      오류 (1단계 OCR 처리): '{item_name_ocr}' 이미지 처리 중 예기치 않은 오류: {e_ocr_general_img}. 건너뜀."
                        if log_func_s1: log_func_s1(err_msg_ocr_gen)
                        logger.error(f"Unexpected error processing image OCR for '{item_name_ocr}': {e_ocr_general_img}", exc_info=True)
                        current_ocr_progress_text = f"'{item_name_ocr}' OCR 오류"


                    # OCR 처리 완료 후 (성공/실패/스킵 무관) 가중치만큼 진행률 업데이트
                    if progress_callback_item_completed and not (stop_event and stop_event.is_set()):
                        progress_callback_item_completed(
                            f"슬라이드 {slide_idx + 1}",    # current_location_info
                            "이미지 처리 완료",          # current_task_type
                            weighted_work_for_ocr_item, # weighted_work_for_item
                            current_ocr_progress_text   # current_text_snippet
                        )
                    if log_func_s1: log_func_s1("\n")

        if stop_event and stop_event.is_set():
            if log_func_s1: log_func_s1(f"--- 1단계: 차트 외 요소 번역 중단됨 ---")
            if f_task_log_s1 and not f_task_log_s1.closed: f_task_log_s1.close()
            return False

        if log_func_s1: log_func_s1(f"--- 1단계: 차트 외 요소 번역 완료 ---\n")
        if f_task_log_s1 and not f_task_log_s1.closed:
            try: f_task_log_s1.close()
            except Exception as e_close_log: logger.warning(f"1단계 핸들러 로그 파일 닫기 실패: {e_close_log}")
        return True